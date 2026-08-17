"""Adapt the Canva Smart Agro deck to the real UCS Meiktila project."""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

SRC = Path(
    r"c:\Users\ASUS\Downloads\Telegram Desktop"
    r"\Smart_Agro_Community_AI_Driven_Farming_Decision_Support_20260803.pptx"
)
OUT = Path(r"d:\SMART-AGRO\docs\Smart-Agro-Community-Presentation.pptx")
OUT_DOWNLOADS = Path(r"c:\Users\ASUS\Downloads\Smart-Agro-Community-Presentation.pptx")

ROOT = Path(r"d:\SMART-AGRO")
IMG_COMMERCIAL = ROOT / "docs" / "Smart-Agro-Commercial.png"
IMG_QR = ROOT / "docs" / "Smart-Agro-Demo-QR.png"
IMG_BLAST = ROOT / "client" / "public" / "demo" / "rice-blast-1.png"
IMG_BLAST2 = ROOT / "client" / "public" / "demo" / "rice-blast-2.png"
IMG_COTTON = ROOT / "client" / "public" / "demo" / "cotton-blight.png"
IMG_LEAFOLDER = ROOT / "client" / "public" / "demo" / "rice-leaffolder.png"
IMG_LOGO = ROOT / "client" / "public" / "logo.png"
IMG_CREST = ROOT / "docs" / "vinyl-sample" / "crest.png"
IMG_TERRACES = ROOT / "client" / "public" / "images" / "footer-terraces.png"

OLIVE = RGBColor(0x47, 0x49, 0x25)
OLIVE_MID = RGBColor(0x70, 0x69, 0x25)
CREAM = RGBColor(0xE8, 0xE4, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2A, 0x2C, 0x14)

# Exact full-shape text (joined with \n) → replacement string or list of paragraphs.
SHAPE_MAP: dict[str, str | list[str]] = {
    "Team Name | University": "UCS Meiktila  ·  2026",
    "Liceria & Co.": "UCS Meiktila",
    "www.reallygreatsite.com": "smart-agro-ucs.surge.sh",
    "www.smartagrocommunity.com": "smart-agro-ucs.surge.sh",
    "AI-Powered Agricultural Decision Support, Disease Detection & GIS Platform  |  MERN Stack + Python AI/ML + GIS": (
        "Smart Tools for Myanmar Farmers  |  Detect · BaGyi Pyoe · Community · Weather · Heatmap"
    ),
    "Farmers across communities face critical barriers that reduce crop yield, increase losses, and limit growth. Without timely information, expert access, and digital tools, these challenges compound — leading to preventable crop failure and economic hardship.": (
        "A yellowing leaf at dawn still means waiting for an officer, guessing a pesticide, or losing a season. English-only apps and distant offices cannot keep up with blast, blight, and pests in Myanmar fields."
    ),
    "Fragmented Information": "English-Only Advice",
    "Limited Expert Access": "Isolated Outbreaks",
    "Difficult Decision Making": "No Trusted Review",
    "Upload crop images for instant AI-powered disease identification with confidence scores and actionable treatment recommendations.": (
        "Photograph a leaf. Gemini names crop and disease across 19 Myanmar crops — with confidence, severity, and DOA treatment."
    ),
    "GIS Agricultural Map": "Outbreak Heatmap",
    "Visualize farm locations, disease hotspots, risk zones, weather markers, and nearby agricultural services on an interactive map.": (
        "Township map of community detections. When neighbours scan the same pest in the same week, the region lights up."
    ),
    "Decision Support": "BaGyi Pyoe",
    "Integrates AI results, weather, and environmental data into a smart engine that assesses risk levels and recommends precise actions.": (
        "A weather-aware farming elder in Myanmar and English. It uses live township weather — it does not invent a forecast."
    ),
    "Connect farmers and experts through Q&A forums, knowledge sharing, and community-driven insights for collaborative smart farming.": (
        "Field posts, photos, linked diagnoses, reports, and expert answers. A farm network with moderation, not noise."
    ),
    "AI-powered image analysis identifies crop diseases at early stages. Farmers upload photos for instant prediction, confidence scoring, and actionable treatment recommendations before losses occur.": (
        "Same-day diagnosis from a field photo: crop, disease, confidence, severity, alternatives, and a bilingual lab report."
    ),
    "Data-Driven Decision Support": "Weather-Aware Advice",
    "Integrates AI results, weather data, and historical records into a smart decision engine. Delivers risk assessments and prioritized recommendations to guide timely, informed farming actions.": (
        "BaGyi Pyoe ties humidity, rain, and crop risk to IPM steps a farmer can take today — in their own language."
    ),
    "GIS Agricultural Intelligence": "Township Heatmap",
    "Location-based mapping visualizes farm zones, disease hotspots, risk overlays, and nearby services. Spatial context empowers farmers with real-time geographic agricultural awareness.": (
        "Leaflet + Myanmar GeoJSON. Detections become early warning at township and region level before the next farm is lost."
    ),
    "Knowledge Sharing & Scalability": "Expert Review & Trust",
    "Connects farmers and experts through a collaborative community platform. Built on a modular MERN + Python AI architecture designed to scale across regions and future smart farming needs.": (
        "Farmers request a second look. Experts verify, correct, or reject with a reason. Models assist; agronomists confirm."
    ),
    "React.js renders the user interface. Node.js + Express.js handles API routing, business logic, and connects all services. Modular REST architecture ensures clean separation of concerns.": (
        "React 19 PWA talks to a typed Express API. Email + OTP, JWT refresh, Zod validation, Helmet, CORS, and rate limits."
    ),
    "MongoDB stores all application data. Python AI Service handles disease detection via machine learning. GIS & External APIs provide spatial intelligence, weather, and location-based data.": (
        "MongoDB + GridFS. Gemini vision and chat with Cursor fallback. Optional FastAPI rice SVM. Open-Meteo. DOA field guides."
    ),
    "Upload crop images for instant AI-powered disease identification. Python ML model analyzes visual symptoms and returns predictions with confidence scores and recommended actions.": (
        "Books, articles, and journals in Myanmar and English. Admin-managed, searchable, and written for the field."
    ),
    "Smart Decision Support": "Expert Review",
    "Data-driven recommendations combining AI results, weather, and environmental context. Helps farmers understand risk levels and take timely, informed action on their crops.": (
        "A queue of real photos and cases. Experts verify or correct the AI result; the farmer receives a notice with the reason."
    ),
    "Interactive map visualizing farm locations, disease hotspots, risk zones, and nearby agricultural services. Spatial intelligence for location-aware farming decisions.": (
        "One-to-one and group messages, friend requests, township profiles, and posts linked to a verified diagnosis."
    ),
    "Real-time and forecast weather data integrated into the platform. Environmental conditions inform disease risk assessments and support proactive crop management decisions.": (
        "Open-Meteo current conditions and 7-day forecast by township or GPS, with crop tips and severe-weather alerts."
    ),
    "Farm Location Mapping": "Township & GPS Pins",
    "Disease Hotspot Heatmap": "Community Heat Layer",
    "Risk Zone Overlays": "Region Choropleth",
    "Weather & Service Markers": "Weather-Linked Risk",
    "GIS-Based": "Outbreak",
    "Agricultural Intelligence.": "Heatmap.",
    "GIS visualization adds critical spatial context to agricultural management. Farmers and agronomists can view real-time farm locations, track disease spread through heatmap overlays, identify high-risk zones, and locate nearby agricultural services — all on one interactive map. Key Stats: Total Farms Monitored | Active Disease Cases | High-Risk Areas Flagged | Nearby Services Found. Spatial intelligence enables faster, location-aware decisions to protect crops and optimize resource deployment.": (
        "Leaflet, Myanmar GeoJSON, and a heat layer of community detections. Filter by disease and date. When many farmers scan the same pest, the township map speaks while the crop can still be saved."
    ),
    "Connecting": "Farmers",
    "Farmers Through Knowledge.": "Learn Together.",
    "React.js Frontend": "React 19 PWA",
    "Python AI Service": "Gemini + FastAPI",
    "Password & Input Security": "OTP & Upload Safety",
    "JWT-based authentication ensures only verified users access the platform. Role-Based Access Control (RBAC) restricts permissions by user type — farmer, expert, or administrator — protecting sensitive data and actions.": (
        "Email + OTP signs the farmer in. JWT access and refresh tokens. Roles: Guest, Farmer, Expert, Admin — experts review cases; admins run the platform."
    ),
    "All passwords are securely hashed using industry-standard algorithms before storage. Strict input validation and sanitization prevent injection attacks, ensuring data integrity across all platform entry points.": (
        "No password store. OTP is rate-limited. Uploads are type-checked by magic bytes. Zod validates every request body before it reaches the database."
    ),
    "API endpoints are protected with rate limiting, token validation, and HTTPS encryption. Unauthorized access attempts are blocked and logged, safeguarding the backend and AI service communication channels.": (
        "Helmet, CORS, and express-rate-limit on the gateway. The Python AI service stays private. Gemini keys rotate when a quota is hit."
    ),
    "Continuous system monitoring detects anomalies and performance issues in real time. Automated database backups and disaster recovery protocols ensure data reliability and platform uptime for all users.": (
        "Admin audit logs, moderation queue, and GridFS file storage. Community reports need a reason; deletions notify the author."
    ),
    "Agriculture": "Smart Agro",
    "Technology.": "Community.",
}

PARA_MAP: dict[tuple[str, ...], list[str]] = {
    (
        "Upload → Preprocessing → Python AI Model → Classification → Result",
        "Crop: Rice / Onion",
        "Prediction: Bacterial Leaf Blight",
        "Confidence: 87%",
        "Risk Level: HIGH",
        "Recommendation: Isolate affected plants immediately. Apply copper-based bactericide. Monitor surrounding crops daily for spread.",
        "Workflow: Farmer uploads crop image → System preprocesses and normalizes → Python AI model analyzes patterns → Disease classified with confidence score → Actionable recommendation delivered instantly.",
        "Disclaimer: AI prediction is decision support only and should be verified by an agricultural professional.",
    ): [
        "Photo → Quality gate → Gemini vision → DOA protocol → Result",
        "Crops: 19 Myanmar crops (rice, onion, chili, cotton…)",
        "Example: Rice blast  ·  Myanmar name first",
        "Confidence + severity (Mild → Critical)",
        "Treatment: IPM first, then DOA chemistry with label-use caution",
        "Farmer can download a bilingual lab report, share to Community, or request expert review.",
        "Workflow: Photograph a leaf → optional township / GPS → AI names crop & disease → guide + report → warn the next farm.",
        "Disclaimer: AI supports the farmer — an expert can still confirm. Not a replacement for a field officer.",
    ],
    (
        "INPUT SIGNALS",
        "Crop Condition",
        "Visual symptoms and growth stage data captured from field observation.",
        "AI Result",
        "Disease prediction with confidence score from image analysis model.",
        "Weather Data",
        "Current temperature, humidity, rainfall, and forecast conditions.",
        "Environmental Data",
        "Soil moisture, historical reports, and surrounding risk indicators.",
        "All inputs converge into the AI Decision Support Engine for unified risk assessment.",
    ): [
        "WHAT THE FARMER BRINGS",
        "Leaf photo",
        "A clear JPEG / PNG / WebP of leaf, stem, fruit, or pest damage (up to 5 MB).",
        "Place",
        "Township search or GPS so the result can feed the outbreak heatmap.",
        "Weather",
        "Live Open-Meteo temperature, humidity, rain, and alerts for that point.",
        "Knowledge",
        "DOA Myanmar field protocols mapped to the detected label.",
        "Photo, place, weather, and DOA guides meet in one result card.",
    ],
    (
        "OUTPUT RECOMMENDATIONS",
        "Risk Level: HIGH",
        "Visible symptoms combined with high humidity and recent outbreak reports.",
        "Risk Factors",
        "Active lesions detected, favorable spread conditions, nearby cases reported.",
        "Recommended Actions",
        "Inspect all plants within 50m radius. Apply appropriate fungicide treatment within 24–48 hours.",
        "Monitoring Advice",
        "Re-assess crop condition every 3 days. Log changes and share findings with community.",
        '"AI prediction supports — not replaces — professional agronomist verification."',
    ): [
        "WHAT THE FARMER GETS",
        "Named problem",
        "Crop and disease in Myanmar first, with confidence and ranked alternatives.",
        "What to do",
        "Symptoms, cultural control, and recommended chemistry — with label-use caution.",
        "Proof",
        "Downloadable bilingual .docx lab report. Optional expert verify / correct / reject.",
        "Warn others",
        "Share to Community. The detection can light up the township heatmap.",
        '"AI supports — it does not replace — a professional agronomist."',
    ],
    (
        '"Leaves are turning yellow. Has anyone experienced this?"',
        "— Farmer, Rice Field Community",
        '💬 Farmer Reply: "Yes, check for nitrogen deficiency or fungal infection."',
        '✅ Community Member: "I had this last season — try foliar spray."',
        '⭐ Agricultural Expert: "Likely early-stage blast disease. Apply recommended fungicide within 48 hours."',
        "👍 Like   💬 Comment   🌿 Helpful   ✔ Expert Answer",
        "Join the community. Ask questions. Share knowledge. Grow together.",
    ): [
        '"Leaves are turning yellow. Has anyone in this township seen this?"',
        "— Farmer, Bago paddy community",
        '💬 Neighbour: "We had this last week — I posted my Detect result."',
        "✅ Linked diagnosis: Rice blast, expert-reviewed.",
        '⭐ Expert: "Scout nearby plots. Follow the DOA blast protocol. Share the heatmap pin."',
        "Like   ·   Comment   ·   Report   ·   Ask an expert",
        "Ask. Share a photo. Warn the next farm — with moderation, not rumour.",
    ],
    (
        "Our platform is built on a modern, modular architecture. React.js powers the frontend, Node.js + Express.js handles the backend logic, MongoDB stores all data, and Python AI runs as a separate service for disease detection. GIS and external APIs enrich the system with spatial intelligence.",
        "Architecture flow:",
        "React.js → Node.js + Express.js → MongoDB",
        "Backend ↔ Python AI Service",
        "Backend ↔ GIS / External APIs",
    ): [
        "A React PWA talks to a typed Express API on MongoDB. Gemini diagnoses photos and powers BaGyi Pyoe. Open-Meteo grounds advice in real township weather. DOA protocols and expert review keep recommendations accountable.",
        "Architecture flow:",
        "React 19 PWA → Express + JWT + Zod → MongoDB",
        "Gateway ↔ Gemini / Cursor  ·  FastAPI rice SVM",
        "Gateway ↔ Open-Meteo  ·  Leaflet heatmap",
    ],
    (
        "✦ Faster access to agricultural information",
        "✦ Earlier disease awareness & detection",
        "✦ Better data-driven decision support",
        "✦ Easier access to agricultural services",
        "✦ Community knowledge sharing & collaboration",
    ): [
        "✦ Same-day Myanmar-language diagnosis from a phone photo",
        "✦ BaGyi Pyoe advice tied to live township weather",
        "✦ Outbreak heatmap so isolation does not hide a regional pest",
        "✦ Expert review and admin moderation for trust",
        "✦ Knowledge Center + community for farmers, women, and youth",
    ],
    (
        "Phase 1: AI Disease Detection",
        "Phase 2: Decision Support + GIS",
        "Phase 3: IoT & Sensor Integration",
        "Phase 4: Predictive Analytics",
        "Phase 5: Large-Scale Smart Farming Ecosystem",
    ): [
        "Near: more DOA crop manuals and saved-report caching",
        "Near: outbreak notices for the farmer’s township",
        "Next: voice input for BaGyi Pyoe (low literacy)",
        "Next: township-office pilots on the expert-review queue",
        "Later: seasonal risk forecasts from heatmap + weather history",
    ],
    (
        "Empowering Farmers with AI, Data and Community.",
        "Detect → Understand → Decide → Act → Monitor",
        "Team: [Team Name] | University: [University Name]",
        "Department: [Department] | Contact: [Email / QR Code]",
    ): [
        "Smart Tools for Myanmar Farmers.",
        "See  ·  Ask  ·  Share  ·  Act",
        "UCS Meiktila  ·  2025–2027 project",
        "Arkar Thet Naing · Khant Zaw · Kaung Myat Tun · Yawai Aung",
    ],
    (
        "AI + DATA + GIS + COMMUNITY = SMARTER FARMING",
    ): [
        "DETECT + BAGYI PYOE + COMMUNITY + HEATMAP = SMARTER FARMING",
    ],
}

# Slide 6 titles that collide with slide 3 if we only use SHAPE_MAP.
SLIDE6_TITLES = {
    "AI Disease Detection": "Knowledge Center",
    "GIS Agricultural Map": "Messages & Friends",
}

SLIDE11_TITLES = {
    "How It": "Product",
    "Works.": "Demo.",
}


def shape_key(shape) -> str:
    if not shape.has_text_frame:
        return ""
    parts = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
    return "\n".join(parts)


def para_key(shape) -> tuple[str, ...]:
    if not shape.has_text_frame:
        return tuple()
    return tuple(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())


def set_run_text(paragraph, text: str) -> None:
    runs = paragraph.runs
    if runs:
        runs[0].text = text
        for run in runs[1:]:
            run.text = ""
    else:
        paragraph.text = text


def set_paragraphs(shape, lines: list[str]) -> None:
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    for i, line in enumerate(lines):
        if i < len(paras):
            set_run_text(paras[i], line)
        else:
            p = tf.add_paragraph()
            p.text = line
            if paras:
                try:
                    p.font.size = paras[0].runs[0].font.size
                    p.font.name = paras[0].runs[0].font.name
                    p.font.color.rgb = paras[0].runs[0].font.color.rgb
                except Exception:
                    pass
    for i in range(len(lines), len(paras)):
        set_run_text(paras[i], "")


def apply_text(shape, slide_idx: int) -> None:
    if not shape.has_text_frame:
        return
    pk = para_key(shape)
    if pk in PARA_MAP:
        set_paragraphs(shape, PARA_MAP[pk])
        return

    key = shape_key(shape)
    if slide_idx == 5 and key in SLIDE6_TITLES:
        set_paragraphs(shape, [SLIDE6_TITLES[key]])
        return
    if slide_idx == 10 and key in SLIDE11_TITLES:
        set_paragraphs(shape, [SLIDE11_TITLES[key]])
        return
    if key in SHAPE_MAP:
        val = SHAPE_MAP[key]
        set_paragraphs(shape, val if isinstance(val, list) else [val])
        return

    # Chrome replacements inside otherwise-untouched boxes.
    for p in shape.text_frame.paragraphs:
        raw = p.text
        if not raw:
            continue
        new = raw
        for old, repl in (
            ("Liceria & Co.", "UCS Meiktila"),
            ("www.reallygreatsite.com", "smart-agro-ucs.surge.sh"),
            ("www.smartagrocommunity.com", "smart-agro-ucs.surge.sh"),
            ("Team Name | University", "UCS Meiktila  ·  2026"),
        ):
            if old in new:
                new = new.replace(old, repl)
        if new != raw:
            set_run_text(p, new)


def walk_shapes(shapes, slide_idx: int) -> None:
    for shape in shapes:
        apply_text(shape, slide_idx)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk_shapes(shape.shapes, slide_idx)


def style_text_frame(tf, size_pt: float, color: RGBColor, bold: bool = False, align=PP_ALIGN.CENTER) -> None:
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size_pt)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Calibri"


def add_label_box(slide, left, top, width, height, title: str, subtitle: str, fill: RGBColor) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    spPr = box._element.spPr
    # Soft corners
    for child in list(spPr):
        if child.tag == qn("a:prstGeom"):
            pass
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(12)
    run2.font.color.rgb = WHITE
    run2.font.name = "Calibri"


def fit_rect(path: Path, left, top, width, height):
    """Keep aspect ratio inside a cell (letterbox)."""
    try:
        from PIL import Image

        with Image.open(path) as im:
            aspect = im.width / float(im.height)
    except Exception:
        return left, top, width, height
    cell_aspect = width / float(height)
    if aspect > cell_aspect:
        new_h = int(width / aspect)
        top = top + (height - new_h) // 2
        height = new_h
    else:
        new_w = int(height * aspect)
        left = left + (width - new_w) // 2
        width = new_w
    return left, top, width, height


def add_picture_or_slot(slide, path: Path | None, left, top, width, height, slot_title: str) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.color.rgb = CREAM
    bg.line.width = Emu(8000)

    cap_h = Emu(280000)
    if path and path.exists():
        pl, pt, pw, ph = fit_rect(path, left, top, width, height - cap_h)
        slide.shapes.add_picture(str(path), pl, pt, pw, ph)
        cap = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top + height - cap_h, width, cap_h
        )
        cap.fill.solid()
        cap.fill.fore_color.rgb = DARK
        cap.line.fill.background()
        tf = cap.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = slot_title
        run.font.size = Pt(10)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"
        return
    add_label_box(slide, left, top, width, height, "PHOTO", slot_title, OLIVE_MID)


def convert_slide_11_to_video(slide) -> None:
    # Clear the long workflow body so the video frame is the focus.
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        key = shape_key(shape)
        if key.startswith("End-to-End Workflow") or key.startswith("① Farmer") or key.startswith(
            "Insert your walkthrough video here"
        ):
            set_paragraphs(shape, [""])

    # 16:9 frame on the right content well (leave left photo group as photo space).
    left, top, width, height = Emu(9000000), Emu(4300000), Emu(7800000), Emu(4387500)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = DARK
    frame.line.color.rgb = CREAM
    frame.line.width = Emu(18000)
    tf = frame.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "▶   DEMO VIDEO"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "PowerPoint: Insert  →  Video  →  drop your MP4 here"
    run2.font.size = Pt(14)
    run2.font.color.rgb = CREAM
    run2.font.name = "Calibri"
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "Leave this frame empty until the clip is ready"
    run3.font.size = Pt(12)
    run3.font.color.rgb = CREAM
    run3.font.name = "Calibri"

    # Photo caption on the existing left image well.
    cap = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(1900000), Emu(7800000), Emu(6200000), Emu(500000)
    )
    cap.fill.solid()
    cap.fill.fore_color.rgb = DARK
    cap.line.fill.background()
    tf = cap.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "PHOTO SPACE  ·  replace with an app screenshot"
    run.font.size = Pt(12)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"


def convert_last_slide_to_gallery(slide) -> None:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        key = shape_key(shape)
        if key.startswith("Smart Tools for Myanmar Farmers") or key.startswith(
            "Empowering Farmers with AI"
        ):
            set_paragraphs(shape, [""])

    # Team names stay readable on the left under Thank You.
    team = slide.shapes.add_textbox(Emu(2775185), Emu(6800000), Emu(6200000), Emu(1400000))
    tf = team.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Arkar Thet Naing  ·  Khant Zaw"
    run.font.size = Pt(14)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "Kaung Myat Tun  ·  Yawai Aung"
    run2.font.size = Pt(14)
    run2.font.color.rgb = WHITE
    run2.font.name = "Calibri"
    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    run3.text = "University of Computer Studies (Meiktila)"
    run3.font.size = Pt(12)
    run3.font.color.rgb = CREAM
    run3.font.name = "Calibri"

    photos = [
        (IMG_COMMERCIAL, "Field scan"),
        (IMG_BLAST, "Rice blast"),
        (IMG_COTTON, "Cotton blight"),
        (IMG_LEAFOLDER, "Rice leaffolder"),
        (IMG_QR, "Live demo QR"),
        (None, "Team photo"),
        (None, "Team photo"),
        (IMG_CREST, "UCS Meiktila"),
    ]

    # 4 columns × 2 rows on the right, sitting over the old copy block.
    origin_x = Emu(9800000)
    origin_y = Emu(3180000)
    cell_w = Emu(1900000)
    cell_h = Emu(2480000)
    gap_x = Emu(140000)
    gap_y = Emu(160000)

    for i, (path, title) in enumerate(photos):
        col = i % 4
        row = i // 4
        left = origin_x + col * (cell_w + gap_x)
        top = origin_y + row * (cell_h + gap_y)
        add_picture_or_slot(slide, path, left, top, cell_w, cell_h, title)


def add_photo_hint_on_detection_slide(slide) -> None:
    """Slide 7 already has a left image well — label it as replaceable photo space."""
    cap = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(2400000), Emu(7800000), Emu(5600000), Emu(480000)
    )
    cap.fill.solid()
    cap.fill.fore_color.rgb = DARK
    cap.line.fill.background()
    tf = cap.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "PHOTO SPACE  ·  drop a Detect screenshot here"
    run.font.size = Pt(12)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"


def main() -> None:
    shutil.copy2(SRC, OUT)
    prs = Presentation(str(OUT))

    for i, slide in enumerate(prs.slides):
        walk_shapes(slide.shapes, i)

    slides = list(prs.slides)
    add_photo_hint_on_detection_slide(slides[6])
    convert_slide_11_to_video(slides[10])
    convert_last_slide_to_gallery(slides[14])

    # Core properties
    prs.core_properties.title = "Smart Agro Community — UCS Meiktila"
    prs.core_properties.author = "University of Computer Studies (Meiktila)"
    prs.core_properties.subject = "Smart Tools for Myanmar Farmers"

    prs.save(str(OUT))
    shutil.copy2(OUT, OUT_DOWNLOADS)
    print(f"Wrote {OUT}")
    print(f"Copied {OUT_DOWNLOADS}")


if __name__ == "__main__":
    main()
