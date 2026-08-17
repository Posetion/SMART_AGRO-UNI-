"""Polish the Smart Agro deck: bullets, alignment, diagrams, roadmap, animations."""

from __future__ import annotations

import shutil
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

PPTX = Path(r"d:\SMART-AGRO\docs\Smart-Agro-Community-Presentation.pptx")
DOWNLOADS = Path(r"c:\Users\ASUS\Downloads\Smart-Agro-Community-Presentation.pptx")
USE_CASE = Path(r"d:\SMART-AGRO\docs\diagrams\01-uml-use-case.png")
SYS_FLOW = Path(r"d:\SMART-AGRO\docs\diagrams\04-system-flow.png")

OLIVE = RGBColor(0x47, 0x49, 0x25)
OLIVE_MID = RGBColor(0x70, 0x69, 0x25)
CREAM = RGBColor(0xE8, 0xE4, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2A, 0x2C, 0x14)


def shape_key(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())


def first_line(shape) -> str:
    if not shape.has_text_frame:
        return ""
    for p in shape.text_frame.paragraphs:
        if p.text.strip():
            return p.text.strip()
    return ""


def set_lines(shape, lines: list[str], align=PP_ALIGN.LEFT) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    paras = list(tf.paragraphs)
    sample = None
    for p in paras:
        if p.runs:
            sample = p.runs[0]
            break
    for i, line in enumerate(lines):
        if i < len(paras):
            p = paras[i]
            if p.runs:
                p.runs[0].text = line
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = line
        else:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
            if sample is not None:
                try:
                    run.font.size = sample.font.size
                    run.font.name = sample.font.name
                    run.font.bold = sample.font.bold
                    run.font.color.rgb = sample.font.color.rgb
                except Exception:
                    pass
        p.alignment = align
        try:
            p.space_after = Pt(4)
        except Exception:
            pass
    for i in range(len(lines), len(paras)):
        p = paras[i]
        if p.runs:
            p.runs[0].text = ""
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.text = ""
        p.alignment = align


def force_left(shape) -> None:
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        if p.alignment == PP_ALIGN.JUSTIFY or p.alignment is None:
            # Keep big titles as they are if already LEFT; convert justify bodies.
            p.alignment = PP_ALIGN.LEFT


def walk(shapes, fn):
    for sh in shapes:
        fn(sh)
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk(sh.shapes, fn)


def bullet(lines: list[str]) -> list[str]:
    out = []
    for ln in lines:
        ln = ln.strip()
        if not ln:
            continue
        if ln.startswith("•") or ln.startswith("–") or ln.startswith("-"):
            out.append(ln)
        else:
            out.append("•  " + ln)
    return out


def polish_text(prs: Presentation) -> None:
    slides = list(prs.slides)

    # Global: justified body → left
    for slide in slides:
        walk(slide.shapes, force_left)

    # --- Slide 2 challenges body ---
    s = slides[1]
    for sh in s.shapes:
        if first_line(sh).startswith("A yellowing leaf"):
            set_lines(sh, bullet([
                "A yellowing leaf still means waiting for an officer",
                "English-only apps miss Myanmar crops and script",
                "Outbreaks stay hidden from neighbouring farms",
                "No trusted expert check before spraying",
            ]))

    # --- Slide 3 solution cards ---
    s = slides[2]
    replacements_3 = {
        "Photograph a leaf.": bullet([
            "Photo of a leaf, stem, or pest damage",
            "19 Myanmar crops — Gemini names the disease",
            "Confidence, severity, and DOA treatment",
        ]),
        "Township map of community": bullet([
            "Township map of community detections",
            "Same pest, same week — the region lights up",
            "Early warning before the next farm is lost",
        ]),
        "A weather-aware farming elder": bullet([
            "Farming advisor in Myanmar and English",
            "Uses live township weather — never invented",
            "IPM steps a farmer can take today",
        ]),
        "Field posts, photos, linked": bullet([
            "Field posts, photos, and linked diagnoses",
            "Reports go to admin — not rumour",
            "Experts can answer on the same card",
        ]),
    }
    apply_startswith(s, replacements_3)

    # --- Slide 4 objectives ---
    s = slides[3]
    apply_startswith(s, {
        "Same-day diagnosis": bullet([
            "Same-day diagnosis from a field photo",
            "Crop, disease, confidence, and alternatives",
            "Bilingual lab report the farmer can keep",
        ]),
        "BaGyi Pyoe ties humidity": bullet([
            "Advice tied to humidity, rain, and crop risk",
            "Spoken in the farmer’s own language",
            "IPM first — no invented brand dosages",
        ]),
        "Leaflet + Myanmar GeoJSON": bullet([
            "Leaflet map + Myanmar GeoJSON boundaries",
            "Detections become township early warning",
            "Filter by disease and date",
        ]),
        "Farmers request a second look": bullet([
            "Farmer can request an expert second look",
            "Verify, correct, or reject — with a reason",
            "Models assist; agronomists confirm",
        ]),
    })

    # --- Slide 5 architecture → use case + system flow ---
    s = slides[4]
    for sh in s.shapes:
        t = first_line(sh)
        if t == "Architecture.":
            set_lines(sh, ["Design."], PP_ALIGN.LEFT)
        elif t == "Frontend & Backend":
            set_lines(sh, ["Use Case"], PP_ALIGN.LEFT)
        elif t.startswith("React 19 PWA talks"):
            set_lines(sh, bullet([
                "Guest, Farmer, Expert, Admin",
                "Public: weather, heatmap, knowledge",
                "JWT: Detect, chat, Community",
            ]))
        elif t == "Data, AI & External Layers":
            set_lines(sh, ["System Flow"], PP_ALIGN.LEFT)
        elif t.startswith("MongoDB + GridFS"):
            set_lines(sh, bullet([
                "Photo → quality gate → Gemini",
                "Weather + township → save diagnosis",
                "Share, expert review, or heatmap pin",
            ]))

    # Cover the two architecture cards with diagrams (skip if already placed)
    has_pics = any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in s.shapes)
    if not has_pics:
        s.shapes.add_picture(str(USE_CASE), Emu(2500000), Emu(3680000), Emu(6200000), Emu(2680000))
        s.shapes.add_picture(str(SYS_FLOW), Emu(9200000), Emu(3680000), Emu(6200000), Emu(2680000))

    # --- Slide 6 features ---
    s = slides[5]
    apply_startswith(s, {
        "Books, articles, and journals": bullet([
            "Books, articles, and journals",
            "Myanmar and English, searchable",
            "Admin-managed for the field",
        ]),
        "A queue of real photos": bullet([
            "Queue of real field photos",
            "Expert verifies or corrects the AI",
            "Farmer gets a notice with the reason",
        ]),
        "One-to-one and group messages": bullet([
            "One-to-one and group messages",
            "Friend requests and township profiles",
            "Posts can link a verified diagnosis",
        ]),
        "Open-Meteo current conditions": bullet([
            "Current conditions and 7-day forecast",
            "Search by township or GPS",
            "Crop tips and severe-weather alerts",
        ]),
    })

    # --- Slide 7 detection ---
    s = slides[6]
    for sh in s.shapes:
        if first_line(sh).startswith("Photo → Quality gate"):
            set_lines(sh, bullet([
                "Photo → quality gate → Gemini → DOA protocol",
                "19 Myanmar crops (rice, onion, chili, cotton…)",
                "Myanmar disease name first, then English",
                "Confidence + severity (Mild → Critical)",
                "IPM first; chemistry with label-use caution",
                "Download a bilingual lab report",
                "Share to Community or request expert review",
                "AI supports the farmer — it does not replace an officer",
            ]))

    # --- Slide 8 From Detection → Who Uses the Platform ---
    s = slides[7]
    for sh in s.shapes:
        t = first_line(sh)
        if t == "From Detection":
            set_lines(sh, ["Who Uses"], PP_ALIGN.LEFT)
        elif t == "to Decision.":
            set_lines(sh, ["the Platform."], PP_ALIGN.LEFT)
        elif t == "WHAT THE FARMER BRINGS":
            set_lines(sh, [
                "GUEST",
                "•  Weather, heatmap, published knowledge",
                "•  No login required",
                "",
                "FARMER",
                "•  Detect disease from a field photo",
                "•  Ask BaGyi Pyoe in Myanmar or English",
                "•  Post, message, report, and warn neighbours",
                "•  Request an expert review",
            ])
        elif t == "WHAT THE FARMER GETS":
            set_lines(sh, [
                "EXPERT",
                "•  Open the diagnosis review queue",
                "•  Verify, correct, or reject the AI label",
                "•  Send a reason back to the farmer",
                "",
                "ADMIN",
                "•  Moderate Community reports",
                "•  Manage users, roles, and knowledge",
                "•  Read audit logs",
            ])

    # --- Slide 9 heatmap body ---
    s = slides[8]
    apply_startswith(s, {
        "Leaflet, Myanmar GeoJSON": bullet([
            "Leaflet + Myanmar GeoJSON boundaries",
            "Heat layer of community detections",
            "Filter by disease and date",
            "When many farmers scan the same pest, the township lights up",
        ]),
    })

    # --- Slide 10 Farmers Learn Together ---
    s = slides[9]
    for sh in s.shapes:
        t = first_line(sh)
        if (
            t.startswith('"Leaves are turning')
            or t.startswith("Ask. Share")
            or t.startswith("•  Post a field")
            or "Bago paddy" in shape_key(sh)
        ):
            set_lines(sh, bullet([
                "Post a field photo — optionally link a diagnosis",
                "Like, comment, and reply in Myanmar",
                "Report spam, harm, or false advice",
                "Admin decides: keep the post or remove it",
                "Expert answers stay on the card",
                "Neighbours see the same pest in the same week",
            ]))
        elif t in ("Farmers", "Community"):
            set_lines(sh, ["Farmers"], PP_ALIGN.LEFT)
        elif t in ("Learn Together.", "that stays safe."):
            set_lines(sh, ["Learn Together."], PP_ALIGN.LEFT)

    # --- Slide 11: video left, photo space right ---
    s = slides[10]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        key = shape_key(sh)
        if "DEMO VIDEO" not in key and "PHOTO SPACE" not in key and "drop your MP4" not in key:
            continue
        if (sh.left or 0) < 8_000_000:
            set_lines(sh, [
                "▶   DEMO VIDEO",
                "Insert → Video → drop your MP4 here",
                "Leave this frame empty until the clip is ready",
            ], PP_ALIGN.CENTER)
        else:
            set_lines(sh, [
                "PHOTO SPACE",
                "Drop a product screenshot here",
            ], PP_ALIGN.CENTER)

    # --- Slide 12 technology stack: remove long body ---
    s = slides[11]
    for sh in s.shapes:
        t = first_line(sh)
        if t.startswith("A React PWA talks") or t.startswith("Architecture flow") or (
            sh.name == "TextBox 14" and not t
        ):
            delete_shape(sh)
            continue
        elif t == "React 19 PWA":
            set_lines(sh, ["React 19 PWA"], PP_ALIGN.LEFT)
        elif t == "Node.js + Express.js":
            set_lines(sh, ["Node.js + Express"], PP_ALIGN.LEFT)
        elif t == "MongoDB Database":
            set_lines(sh, ["MongoDB + GridFS"], PP_ALIGN.LEFT)
        elif t == "Gemini + FastAPI":
            set_lines(sh, ["Gemini + FastAPI"], PP_ALIGN.LEFT)

    # --- Slide 13 security bullets + real auth ---
    s = slides[12]
    apply_startswith(s, {
        "Email + OTP signs": bullet([
            "Email + password, or guest session",
            "JWT access token + refresh token",
            "Roles: Guest · Farmer · Expert · Admin",
        ]),
        "No password store": bullet([
            "Passwords hashed before storage",
            "Uploads checked by file magic bytes",
            "Zod validates every request body",
        ]),
        "Helmet, CORS, and express": bullet([
            "Helmet, CORS, and rate limits",
            "Python AI service stays private",
            "Gemini keys rotate when a quota is hit",
        ]),
        "Admin audit logs": bullet([
            "Admin audit logs and moderation queue",
            "Reports need a reason",
            "Deletions notify the author",
        ]),
    })
    for sh in s.shapes:
        if first_line(sh) == "OTP & Upload Safety":
            set_lines(sh, ["Password & Upload Safety"], PP_ALIGN.LEFT)

    # --- Slide 14 Future Roadmap only ---
    rework_roadmap(slides[13])


def apply_startswith(slide, mapping: dict[str, list[str]]) -> None:
    for sh in slide.shapes:
        t = first_line(sh)
        for prefix, lines in mapping.items():
            if t.startswith(prefix):
                set_lines(sh, lines)
                break


def delete_shape(shape) -> None:
    el = shape._element
    el.getparent().remove(el)


def rework_roadmap(slide) -> None:
    # Drop the old two-column Impact / Vision layout entirely.
    to_remove = []
    for sh in slide.shapes:
        t = first_line(sh)
        name = (sh.name or "")
        if t in ("Impact &", "Future"):
            set_lines(sh, ["Future"], PP_ALIGN.LEFT)
            continue
        if t in ("Future Vision.", "Roadmap."):
            set_lines(sh, ["Roadmap."], PP_ALIGN.LEFT)
            continue
        generated = (
            name.startswith("Rounded Rectangle")
            or name.startswith("Rectangle")
            or name.startswith("Oval")
        )
        if generated and (sh.top or 0) > 3_000_000:
            to_remove.append(sh)
            continue
        if t in ("Current Impact", "Future Roadmap"):
            to_remove.append(sh)
        elif t.startswith("✦ Same-day") or t.startswith("Near:") or t.startswith("01"):
            to_remove.append(sh)
        elif name.startswith("Group ") or (name.startswith("Freeform ") and sh.top and sh.top > 3_000_000):
            if sh.shape_id >= 12:
                to_remove.append(sh)
        elif name.startswith("TextBox ") and sh.shape_id in (13, 14, 16, 17):
            to_remove.append(sh)
    for sh in to_remove:
        try:
            delete_shape(sh)
        except Exception:
            pass

    # Unique timeline band — not the 2-column card grid used elsewhere
    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(2100000), Emu(3280000), Emu(14200000), Emu(5080000)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0xF4, 0xF1, 0xE6)
    band.line.fill.background()

    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(2800000), Emu(3920000), Emu(12800000), Emu(50000)
    )
    rail.fill.solid()
    rail.fill.fore_color.rgb = OLIVE_MID
    rail.line.fill.background()

    phases = [
        ("01", "NEAR", "More DOA crop manuals in the treatment catalog"),
        ("02", "NEAR", "Cache saved reports for weak connections"),
        ("03", "NEAR", "Outbreak notices for the farmer’s township"),
        ("04", "NEXT", "Voice input for BaGyi Pyoe (low literacy)"),
        ("05", "LATER", "Seasonal risk forecasts from heatmap + weather"),
    ]
    origin_x = Emu(2480000)
    card_w = Emu(2520000)
    gap = Emu(160000)
    top = Emu(4180000)
    for i, (num, when, text) in enumerate(phases):
        left = origin_x + i * (card_w + gap)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left + Emu(980000), Emu(3760000), Emu(360000), Emu(360000)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = DARK if i % 2 == 0 else OLIVE
        dot.line.fill.background()

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, Emu(3920000)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK if i % 2 == 0 else OLIVE
        card.line.fill.background()
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(140000)
        tf.margin_right = Emu(140000)
        tf.margin_top = Emu(180000)
        tf.margin_bottom = Emu(140000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = num
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = CREAM
        run.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = when
        r2.font.size = Pt(11)
        r2.font.bold = True
        r2.font.color.rgb = CREAM
        r2.font.name = "Calibri"
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        r3 = p3.add_run()
        r3.text = text
        r3.font.size = Pt(13)
        r3.font.color.rgb = WHITE
        r3.font.name = "Calibri"


def add_timing_xml(slide, shape_ids: list[int]) -> None:
    """Entrance fade, one click per shape — injected as p:timing OOXML."""
    # Build a simple sequence: each shape fades in on click.
    nsmap = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }

    def ptag(name):
        return "{http://schemas.openxmlformats.org/presentationml/2006/main}" + name

    cSld = slide._element.cSld
    # Remove existing timing
    for child in list(slide._element):
        if child.tag == ptag("timing"):
            slide._element.remove(child)

    timing = etree.SubElement(slide._element, ptag("timing"))
    tnLst = etree.SubElement(timing, ptag("tnLst"))
    par_root = etree.SubElement(tnLst, ptag("par"))
    cTn_root = etree.SubElement(par_root, ptag("cTn"), {
        "id": "1", "dur": "indefinite", "restart": "never", "nodeType": "tmRoot",
    })
    child_root = etree.SubElement(cTn_root, ptag("childTnLst"))
    seq = etree.SubElement(child_root, ptag("seq"), {"concurrent": "true", "nextAc": "seek"})
    cTn_seq = etree.SubElement(seq, ptag("cTn"), {
        "id": "2", "dur": "indefinite", "nodeType": "mainSeq",
    })
    child_seq = etree.SubElement(cTn_seq, ptag("childTnLst"))

    nid = 3
    for i, spid in enumerate(shape_ids):
        par = etree.SubElement(child_seq, ptag("par"))
        cTn = etree.SubElement(par, ptag("cTn"), {
            "id": str(nid), "fill": "hold",
        })
        nid += 1
        stCondLst = etree.SubElement(cTn, ptag("stCondLst"))
        cond = etree.SubElement(stCondLst, ptag("cond"), {
            "delay": "0" if i else "0",
        })
        if i == 0:
            cond.set("delay", "0")
        else:
            # on click
            pass
        # click trigger via cond/cTn on previous — use onClick via tgtEl
        # Simpler: first with previous (after slide load), rest on click
        if i == 0:
            cond.set("delay", "0")
        else:
            cond.set("evt", "onClick")
            cond.set("delay", "0")
            tgtEl = etree.SubElement(cond, ptag("tgtEl"))
            etree.SubElement(tgtEl, ptag("sldTgt"))

        child = etree.SubElement(cTn, ptag("childTnLst"))
        par2 = etree.SubElement(child, ptag("par"))
        cTn2 = etree.SubElement(par2, ptag("cTn"), {
            "id": str(nid), "fill": "hold",
        })
        nid += 1
        st2 = etree.SubElement(cTn2, ptag("stCondLst"))
        etree.SubElement(st2, ptag("cond"), {"delay": "0"})
        child2 = etree.SubElement(cTn2, ptag("childTnLst"))

        # effect: fade
        animEffect = etree.SubElement(child2, ptag("animEffect"), {
            "transition": "in", "filter": "fade",
        })
        cBhvr = etree.SubElement(animEffect, ptag("cBhvr"))
        cTnE = etree.SubElement(cBhvr, ptag("cTn"), {
            "id": str(nid), "dur": "500",
        })
        nid += 1
        tgt = etree.SubElement(cBhvr, ptag("tgtEl"))
        etree.SubElement(tgt, ptag("spTgt"), {"spid": str(spid)})

    prevCondLst = etree.SubElement(seq, ptag("prevCondLst"))
    condp = etree.SubElement(prevCondLst, ptag("cond"), {"evt": "onPrev", "delay": "0"})
    tgtp = etree.SubElement(condp, ptag("tgtEl"))
    etree.SubElement(tgtp, ptag("sldTgt"))
    nextCondLst = etree.SubElement(seq, ptag("nextCondLst"))
    condn = etree.SubElement(nextCondLst, ptag("cond"), {"evt": "onNext", "delay": "0"})
    tgtn = etree.SubElement(condn, ptag("tgtEl"))
    etree.SubElement(tgtn, ptag("sldTgt"))


def collect_anim_ids(slide) -> list[int]:
    ids = []
    for sh in slide.shapes:
        name = (sh.name or "").lower()
        if "freeform" in name:
            continue
        if sh.has_text_frame and shape_key(sh):
            # skip tiny chrome
            t = first_line(sh)
            if t.startswith("Page ") or t in ("Smart Agro", "Community.", "UCS Meiktila", "smart-agro-ucs.surge.sh", "Smart Agro Community"):
                continue
            ids.append(sh.shape_id)
        else:
            # pictures / cards we added
            if sh.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.AUTO_SHAPE):
                if sh.has_text_frame and not shape_key(sh):
                    continue
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    ids.append(sh.shape_id)
                elif first_line(sh) or (sh.name or "").startswith("Rounded"):
                    ids.append(sh.shape_id)
    return ids


def main() -> None:
    prs = Presentation(str(PPTX))
    polish_text(prs)
    prs.save(str(PPTX))
    shutil.copy2(PPTX, DOWNLOADS)
    print("saved text polish", PPTX)


if __name__ == "__main__":
    main()
