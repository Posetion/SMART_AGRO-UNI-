"""List every shape with position so we can hide leftovers."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = r"d:\SMART-AGRO\docs\Smart-Agro-Community-Presentation.pptx"
prs = Presentation(PPTX)

def walk(shapes, indent=0):
    for sh in shapes:
        kind = str(sh.shape_type).split(".")[-1] if sh.shape_type else "?"
        name = sh.name or ""
        txt = ""
        if sh.has_text_frame:
            lines = [p.text.strip() for p in sh.text_frame.paragraphs if p.text.strip()]
            txt = " | ".join(lines)[:140]
        print(
            f"{'  '*indent}id={sh.shape_id} [{kind}] {name} "
            f"@({sh.left},{sh.top}) {sh.width}x{sh.height} {txt}"
        )
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk(sh.shapes, indent + 1)

for i in (5, 10, 11, 12, 14):
    print(f"\n===== SLIDE {i} ALL SHAPES =====")
    walk(prs.slides[i - 1].shapes)
