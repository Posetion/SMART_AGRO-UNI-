"""Dump slide titles and body text for a quick visual check."""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = r"d:\SMART-AGRO\docs\Smart-Agro-Community-Presentation.pptx"
prs = Presentation(PPTX)

def walk(shapes, indent=0):
    for sh in shapes:
        kind = str(sh.shape_type).split(".")[-1] if sh.shape_type else "?"
        name = sh.name or ""
        txt = ""
        if sh.has_text_frame:
            lines = [p.text.strip() for p in sh.text_frame.paragraphs if p.text.strip()]
            txt = " | ".join(lines)[:220]
        if txt or sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print(f"{'  '*indent}[{kind}] {name}: {txt or '(picture)'}")
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk(sh.shapes, indent + 1)

for i, slide in enumerate(prs.slides, 1):
    print(f"\n===== SLIDE {i} =====")
    walk(slide.shapes)
